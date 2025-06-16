from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(title, scope, literature_review, methodology, milestones, support_needed):
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = title

    # Project Scope Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Project Scope"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = scope

    # Literature Review / Market Study Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Literature Review / Market Study"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = literature_review

    # Methodology Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Methodology"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = methodology

    # Planned Achievements Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Planned Achievements"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = milestones

    # Support Needed Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Support Needed from Senzmate"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = support_needed

    prs.save("project_proposal.pptx")

if __name__ == "__main__":
    title = "Real-Time Sentiment Analysis of Social Media Feeds"
    scope = "Our project aims to develop a real-time sentiment analysis system for monitoring social media feeds. This system will analyze public sentiment and brand perception across various social media platforms."
    literature_review = "Research indicates a growing importance of sentiment analysis in understanding consumer behavior and brand reputation. Market trends show an increasing demand for real-time analytics to drive proactive decision-making."
    methodology = "1. Data Acquisition: Utilize Twitter API for real-time data ingestion.\n2. Data Processing: Implement Spark for streaming data processing and sentiment analysis using NLP techniques.\n3. Data Storage: Store analyzed data in Elasticsearch for efficient querying and indexing.\n4. Visualization: Use Kibana for real-time visualization of sentiment trends and insights."
    milestones = "- Milestone 1 (Month 2): Complete setup of data ingestion and initial Spark processing pipeline.\n- Milestone 2 (Month 4): Achieve real-time sentiment analysis and integration with Elasticsearch.\n- Final Outcome: Deliver a functional system with a user-friendly Kibana dashboard for visualizing sentiment trends and insights."
    support_needed = "1. Technical Support: Assistance with configuring and optimizing data ingestion through APIs.\n2. Infrastructure: Access to scalable resources for Spark processing and Elasticsearch storage.\n3. Training: Optional training sessions on advanced Elasticsearch querying and Kibana dashboard customization."

    create_presentation(title, scope, literature_review, methodology, milestones, support_needed)
